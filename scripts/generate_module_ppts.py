#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
모듈별 PPT 생성 스크립트
하나의 PPT 파일에 모든 모듈을 포함합니다. 각 모듈은 하나의 슬라이드로 표현됩니다.
"""

import json
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
# MSO_TABLE_STYLE은 일부 버전에서 사용 불가능할 수 있으므로 조건부 import
try:
    from pptx.enum.table import MSO_TABLE_STYLE
except ImportError:
    # MSO_TABLE_STYLE이 없으면 None으로 설정 (사용하지 않음)
    MSO_TABLE_STYLE = None
from openai import OpenAI
import re
import platform
import matplotlib
matplotlib.use('Agg')  # GUI 백엔드 없이 사용
# 한글 폰트 경고 무시
import warnings
warnings.filterwarnings('ignore', category=UserWarning, message='.*Glyph.*missing from font.*')
import matplotlib.pyplot as plt
import numpy as np
from io import BytesIO
import tempfile

# OpenAI API 키 설정 (환경 변수에서 가져오기)
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")

if OPENAI_API_KEY:
    client = OpenAI(api_key=OPENAI_API_KEY)
else:
    client = None
    print("경고: OPENAI_API_KEY가 설정되지 않았습니다. AI 설명을 생성할 수 없습니다.")

# Gemini API 키 설정 (환경 변수에서 가져오기)
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY") or os.getenv("API_KEY")

# Shape 타입 제외 목록
EXCLUDED_TYPES = ['TextBox', 'GroupBox']

def get_download_folder():
    """사용자의 다운로드 폴더 경로 반환"""
    if platform.system() == "Windows":
        import winreg
        sub_key = r'SOFTWARE\Microsoft\Windows\CurrentVersion\Explorer\Shell Folders'
        downloads_guid = '{374DE290-123F-4565-9164-39C4925E467B}'
        with winreg.OpenKey(winreg.HKEY_CURRENT_USER, sub_key) as key:
            location = winreg.QueryValueEx(key, downloads_guid)[0]
        return location
    elif platform.system() == "Darwin":  # macOS
        return os.path.join(os.path.expanduser("~"), "Downloads")
    else:  # Linux
        return os.path.join(os.path.expanduser("~"), "Downloads")

def sanitize_filename(name):
    """파일명으로 사용할 수 없는 문자 제거"""
    name = re.sub(r'[<>:"/\\|?*]', '_', name)
    name = name.replace(' ', '_')
    return name

def get_input_data_info(module, all_modules, connections):
    """모듈의 입력 데이터 정보 추출"""
    inputs = module.get('inputs', [])
    if not inputs:
        return "입력 데이터 없음"
    
    input_info = []
    for input_port in inputs:
        # 연결된 모듈 찾기
        connected_module = None
        for conn in connections:
            if conn.get('to', {}).get('moduleId') == module.get('id') and \
               conn.get('to', {}).get('portName') == input_port.get('name'):
                from_module_id = conn.get('from', {}).get('moduleId')
                connected_module = next((m for m in all_modules if m.get('id') == from_module_id), None)
                break
        
        port_type = input_port.get('type', 'unknown')
        if connected_module:
            module_name = connected_module.get('name', connected_module.get('type'))
            output_data = connected_module.get('outputData')
            
            if output_data and output_data.get('type') == 'DataPreview':
                columns = output_data.get('columns', [])
                row_count = output_data.get('totalRowCount', 0)
                input_info.append(f"• {input_port['name']} ({port_type}): {module_name}에서 전달\n  - 형태: {row_count}행 × {len(columns)}열")
                if columns:
                    col_names = [col.get('name', '') for col in columns[:5]]
                    input_info.append(f"  - 주요 컬럼: {', '.join(col_names)}{'...' if len(columns) > 5 else ''}")
            else:
                input_info.append(f"• {input_port['name']} ({port_type}): {module_name}에서 전달")
        else:
            input_info.append(f"• {input_port['name']} ({port_type}): 연결되지 않음")
    
    return '\n'.join(input_info) if input_info else "입력 데이터 없음"

def get_output_data_info(module):
    """모듈의 출력 데이터 정보 추출"""
    output_data = module.get('outputData')
    outputs = module.get('outputs', [])
    
    if not outputs:
        return "출력 데이터 없음"
    
    output_info = []
    for output_port in outputs:
        port_type = output_port.get('type', 'unknown')
        port_name = output_port.get('name', 'unknown')
        
        if output_data:
            if output_data.get('type') == 'DataPreview':
                columns = output_data.get('columns', [])
                row_count = output_data.get('totalRowCount', 0)
                output_info.append(f"• {port_name} ({port_type}): 데이터 테이블\n  - 형태: {row_count}행 × {len(columns)}열")
                if columns:
                    col_names = [col.get('name', '') for col in columns[:5]]
                    output_info.append(f"  - 주요 컬럼: {', '.join(col_names)}{'...' if len(columns) > 5 else ''}")
            elif output_data.get('type') == 'StatisticsOutput':
                output_info.append(f"• {port_name} ({port_type}): 통계 분석 결과")
            elif output_data.get('type') == 'TrainedModelOutput':
                output_info.append(f"• {port_name} ({port_type}): 훈련된 모델")
            elif output_data.get('type') == 'StatsModelsResultOutput':
                output_info.append(f"• {port_name} ({port_type}): 통계 모델 결과")
            elif output_data.get('type') == 'EvaluationOutput':
                output_info.append(f"• {port_name} ({port_type}): 모델 평가 결과")
            else:
                output_info.append(f"• {port_name} ({port_type}): {output_data.get('type', '알 수 없는 타입')}")
        else:
            output_info.append(f"• {port_name} ({port_type}): 실행되지 않음")
    
    return '\n'.join(output_info) if output_info else "출력 데이터 없음"

def get_model_equation(output_data):
    """ResultModel의 함수식 생성"""
    if not output_data or output_data.get('type') != 'StatsModelsResultOutput':
        return None
    
    model_type = output_data.get('modelType', 'Unknown')
    summary = output_data.get('summary', {})
    coefficients = summary.get('coefficients', {})
    feature_columns = output_data.get('featureColumns', [])
    label_column = output_data.get('labelColumn', 'y')
    
    if not coefficients:
        return None
    
    # 절편(intercept) 찾기
    intercept = coefficients.get('const', {}).get('coef', 0)
    
    # 함수식 생성
    equation_parts = []
    for feature in feature_columns:
        coef_info = coefficients.get(feature, {})
        coef = coef_info.get('coef', 0)
        if coef != 0:
            if coef > 0:
                equation_parts.append(f"+ {coef:.4f}×{feature}")
            else:
                equation_parts.append(f"{coef:.4f}×{feature}")
    
    equation = f"{label_column} = {intercept:.4f}"
    if equation_parts:
        equation += " " + " ".join(equation_parts)
    
    return equation

def generate_insight_headline(module, all_modules=None, connections=None):
    """AI를 사용하여 분석적 통찰이 담긴 헤드라인 생성"""
    module_type = module.get('type', 'Unknown')
    module_name = module.get('name', module_type)
    parameters = module.get('parameters', {})
    output_data = module.get('outputData')
    
    # 기본 헤드라인 (AI가 없을 경우)
    default_headline = f"{module_name} 분석 결과"
    
    # Gemini API 사용 시도
    if GEMINI_API_KEY:
        try:
            try:
                import google.generativeai as genai
            except ImportError:
                # google-generativeai 패키지가 없으면 건너뛰기
                # print("google-generativeai 패키지가 설치되지 않았습니다. 기본 헤드라인을 사용합니다.")
                return default_headline
            except Exception as e:
                # 기타 import 오류도 처리
                # print(f"google-generativeai import 오류: {e}. 기본 헤드라인을 사용합니다.")
                return default_headline
            
            try:
                genai.configure(api_key=GEMINI_API_KEY)
                model = genai.GenerativeModel('gemini-2.5-flash')
            except Exception as e:
                # configure 또는 모델 생성 오류
                # print(f"Gemini API 설정 오류: {e}. 기본 헤드라인을 사용합니다.")
                return default_headline
            
            # 모듈 정보 수집
            module_info = {
                'type': module_type,
                'name': module_name,
                'parameters': parameters
            }
            
            # 입력/출력 데이터 요약
            input_summary = get_input_data_info(module, all_modules, connections)
            output_summary = get_output_data_info(module)
            
            # 분석 내용
            analysis_desc = get_analysis_description(module, all_modules, connections)
            
            # 전체 프로젝트 컨텍스트 수집
            project_context = ""
            if all_modules:
                module_names = [m.get('name', m.get('type', 'Unknown')) for m in all_modules if m.get('type') not in EXCLUDED_TYPES]
                if module_names:
                    project_context = f"이 프로젝트에는 총 {len(module_names)}개의 모듈이 포함되어 있습니다: {', '.join(module_names[:10])}{'...' if len(module_names) > 10 else ''}"
            
            # 이전 모듈 정보 수집 (연결된 입력 모듈)
            previous_modules_info = ""
            if connections:
                input_modules = []
                for conn in connections:
                    if conn.get('to', {}).get('moduleId') == module.get('id'):
                        from_id = conn.get('from', {}).get('moduleId')
                        from_module = next((m for m in all_modules if m.get('id') == from_id), None)
                        if from_module:
                            input_modules.append(from_module.get('name', from_module.get('type', 'Unknown')))
                if input_modules:
                    previous_modules_info = f"이 모듈은 다음 모듈들의 결과를 입력으로 받습니다: {', '.join(input_modules)}"
            
            prompt = f"""당신은 데이터 분석 전문가입니다. 다음 모듈의 분석 결과를 바탕으로, 단순한 모듈 이름이 아닌 **분석적 통찰이 담긴 헤드라인**을 작성해주세요.

[프로젝트 전체 컨텍스트]
{project_context}

[이전 모듈 정보]
{previous_modules_info if previous_modules_info else "이 모듈은 프로젝트의 시작점입니다."}

[모듈 정보]
- 모듈 타입: {module_type}
- 모듈 이름: {module_name}
- 주요 파라미터: {json.dumps(parameters, ensure_ascii=False)[:200]}

[입력 데이터]
{input_summary[:300]}

[출력 데이터]
{output_summary[:300]}

[분석 내용]
{analysis_desc[:500]}

[요구사항 - 반드시 준수하세요]
1. **핵심 원칙**: 단순한 모듈 이름(예: "Scaling Transform 결과", "K-Means 클러스터링")을 나열하지 마세요. 대신 **구체적인 분석 결과나 인사이트**를 담은 헤드라인을 작성하세요.

2. **좋은 예시**:
   - "데이터 스케일링을 통해 변수 간 단위 차이를 제거하고 모델의 수렴 속도를 20% 향상함"
   - "K-Means 클러스터링으로 3개의 고유한 고객 세그먼트를 발견하고 각 그룹의 평균 구매액이 15% 차이를 보임"
   - "결측값을 평균값으로 대체하여 데이터 손실 없이 1,234개 샘플을 완전한 데이터셋으로 변환함"
   - "선형 회귀 모델이 R² 0.85로 주택 가격 예측에 효과적이며, 방 개수와 위치가 가장 중요한 변수로 확인됨"

3. **나쁜 예시 (피해야 할 것)**:
   - "Scaling Transform 결과"
   - "K-Means 클러스터링 완료"
   - "데이터 전처리 수행"
   - "모델 훈련 결과"

4. **작성 가이드**:
   - 한국어로 작성하세요.
   - 60-70자 이내로 간결하게 작성하세요 (50자보다 여유 있게).
   - 숫자, 비율, 구체적인 변화가 있다면 반드시 포함하세요 (예: "20% 향상", "3개 그룹", "1,234개 샘플").
   - 모듈이 수행한 작업의 **결과나 효과**를 강조하세요.
   - 가능하면 정량적 지표나 구체적인 변화를 포함하세요.

5. **출력 형식**:
   - 헤드라인만 작성하고, 다른 설명이나 부가 텍스트는 추가하지 마세요.
   - 따옴표나 인용 부호 없이 순수한 텍스트만 작성하세요.
   - 마침표로 끝나지 않도록 하세요 (선택사항).

헤드라인만 작성하고, 다른 설명은 추가하지 마세요:"""
            
            response = model.generate_content(prompt)
            headline = response.text.strip()
            
            # 불필요한 따옴표 제거
            headline = headline.strip('"\'')
            
            if headline and len(headline) > 5:
                return headline
        except Exception as e:
            # Gemini API 오류는 무시하고 기본 헤드라인 사용
            # print(f"Gemini API 헤드라인 생성 실패: {e}")
            pass
    
    # OpenAI API 사용 시도
    if client:
        try:
            # 전체 프로젝트 컨텍스트 수집
            project_context = ""
            if all_modules:
                module_names = [m.get('name', m.get('type', 'Unknown')) for m in all_modules if m.get('type') not in EXCLUDED_TYPES]
                if module_names:
                    project_context = f"이 프로젝트에는 총 {len(module_names)}개의 모듈이 포함되어 있습니다: {', '.join(module_names[:10])}{'...' if len(module_names) > 10 else ''}"
            
            # 이전 모듈 정보 수집 (연결된 입력 모듈)
            previous_modules_info = ""
            if connections:
                input_modules = []
                for conn in connections:
                    if conn.get('to', {}).get('moduleId') == module.get('id'):
                        from_id = conn.get('from', {}).get('moduleId')
                        from_module = next((m for m in all_modules if m.get('id') == from_id), None)
                        if from_module:
                            input_modules.append(from_module.get('name', from_module.get('type', 'Unknown')))
                if input_modules:
                    previous_modules_info = f"이 모듈은 다음 모듈들의 결과를 입력으로 받습니다: {', '.join(input_modules)}"
            
            prompt = f"""당신은 데이터 분석 전문가입니다. 다음 모듈의 분석 결과를 바탕으로, 단순한 모듈 이름이 아닌 **분석적 통찰이 담긴 헤드라인**을 작성해주세요.

프로젝트 전체 컨텍스트: {project_context}

이전 모듈 정보: {previous_modules_info if previous_modules_info else "이 모듈은 프로젝트의 시작점입니다."}

모듈 타입: {module_type}
모듈 이름: {module_name}
주요 파라미터: {json.dumps(parameters, ensure_ascii=False)[:200]}

입력 데이터: {get_input_data_info(module, all_modules, connections)[:300]}
출력 데이터: {get_output_data_info(module)[:300]}
분석 내용: {get_analysis_description(module, all_modules, connections)[:500]}

요구사항 - 반드시 준수하세요:
1. 핵심 원칙: 단순한 모듈 이름(예: "Scaling Transform 결과", "K-Means 클러스터링")을 나열하지 마세요. 대신 구체적인 분석 결과나 인사이트를 담은 헤드라인을 작성하세요.

2. 좋은 예시:
   - "데이터 스케일링을 통해 변수 간 단위 차이를 제거하고 모델의 수렴 속도를 20% 향상함"
   - "K-Means 클러스터링으로 3개의 고유한 고객 세그먼트를 발견하고 각 그룹의 평균 구매액이 15% 차이를 보임"
   - "결측값을 평균값으로 대체하여 데이터 손실 없이 1,234개 샘플을 완전한 데이터셋으로 변환함"
   - "선형 회귀 모델이 R² 0.85로 주택 가격 예측에 효과적이며, 방 개수와 위치가 가장 중요한 변수로 확인됨"

3. 나쁜 예시 (피해야 할 것):
   - "Scaling Transform 결과"
   - "K-Means 클러스터링 완료"
   - "데이터 전처리 수행"
   - "모델 훈련 결과"

4. 작성 가이드:
   - 한국어로 작성하세요.
   - 60-70자 이내로 간결하게 작성하세요 (50자보다 여유 있게).
   - 숫자, 비율, 구체적인 변화가 있다면 반드시 포함하세요 (예: "20% 향상", "3개 그룹", "1,234개 샘플").
   - 모듈이 수행한 작업의 결과나 효과를 강조하세요.
   - 가능하면 정량적 지표나 구체적인 변화를 포함하세요.

5. 출력 형식:
   - 헤드라인만 작성하고, 다른 설명이나 부가 텍스트는 추가하지 마세요.
   - 따옴표나 인용 부호 없이 순수한 텍스트만 작성하세요.
   - 마침표로 끝나지 않도록 하세요 (선택사항).

헤드라인만 작성하세요:"""
            
            response = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                    {"role": "system", "content": "당신은 데이터 분석 전문가입니다."},
                    {"role": "user", "content": prompt}
                ],
                max_tokens=100,
                temperature=0.7
            )
            
            headline = response.choices[0].message.content.strip()
            headline = headline.strip('"\'')
            
            if headline and len(headline) > 5:
                return headline
        except Exception as e:
            print(f"OpenAI API 헤드라인 생성 실패: {e}")
    
    return default_headline

def create_data_summary_table(slide, module, all_modules, connections, x, y, width, height):
    """입력/출력 데이터 요약 표 생성 (입력 데이터 오른쪽에 속성값, Train Model 출력 데이터 오른쪽에 파라미터 결과)"""
    try:
        # 입력 데이터 정보 수집
        input_info = get_input_data_info(module, all_modules, connections)
        output_info = get_output_data_info(module)
        module_type = module.get('type', 'Unknown')
        parameters = module.get('parameters', {})
        output_data = module.get('outputData')
        
        # 모듈 속성값(parameters) 포맷팅
        def format_parameters(params):
            """파라미터를 읽기 쉬운 문자열로 변환"""
            if not params:
                return "없음"
            
            param_items = []
            for key, value in list(params.items())[:5]:  # 최대 5개만 표시
                if isinstance(value, (dict, list)):
                    value_str = json.dumps(value, ensure_ascii=False)
                    if len(value_str) > 30:
                        value_str = value_str[:27] + "..."
                else:
                    value_str = str(value)
                    if len(value_str) > 30:
                        value_str = value_str[:27] + "..."
                param_items.append(f"{key}: {value_str}")
            
            result = ", ".join(param_items)
            if len(params) > 5:
                result += f" ... (총 {len(params)}개)"
            return result
        
        # Train Model의 파라미터 결과 포맷팅
        def format_train_model_results(output_data):
            """Train Model의 출력 데이터에서 파라미터 결과 포맷팅"""
            if not output_data or output_data.get('type') != 'TrainedModelOutput':
                return None
            
            result_parts = []
            
            # 모델 타입
            model_type = output_data.get('modelType', 'Unknown')
            result_parts.append(f"모델: {model_type}")
            
            # 주요 지표 (metrics)
            metrics = output_data.get('metrics', {})
            if metrics:
                key_metrics = []
                # 주요 지표만 선택
                for key in ['Accuracy', 'R-squared', 'R²', 'RMSE', 'MAE', 'F1-Score', 'Precision', 'Recall']:
                    if key in metrics:
                        value = metrics[key]
                        if isinstance(value, (int, float)):
                            key_metrics.append(f"{key}: {value:.4f}")
                        else:
                            key_metrics.append(f"{key}: {value}")
                
                if key_metrics:
                    result_parts.append(", ".join(key_metrics[:3]))  # 최대 3개 지표
            
            # 계수 정보 (간단히)
            coefficients = output_data.get('coefficients', {})
            intercept = output_data.get('intercept', 0)
            if coefficients:
                coeff_count = len(coefficients)
                result_parts.append(f"계수: {coeff_count}개, 절편: {intercept:.4f}")
            
            return " | ".join(result_parts) if result_parts else None
        
        # 표 데이터 준비 (3열: 항목, 내용, 추가 정보)
        table_data = [
            ['항목', '내용', '추가 정보'],
            ['입력 데이터', input_info.split('\n')[0] if input_info else '없음', format_parameters(parameters)],
        ]
        
        # 출력 데이터 행 (Train Model인 경우 파라미터 결과 추가)
        output_content = output_info.split('\n')[0] if output_info else '없음'
        if module_type == 'TrainModel' and output_data:
            train_results = format_train_model_results(output_data)
            output_additional = train_results if train_results else "없음"
        else:
            output_additional = "없음"
        
        table_data.append(['출력 데이터', output_content, output_additional])
        
        # 입력 데이터 상세 정보 추가
        if '행 ×' in input_info:
            for line in input_info.split('\n'):
                if '행 ×' in line:
                    table_data.append(['입력 형태', line.split(':')[1].strip() if ':' in line else line.strip(), ''])
                    break
        
        # 출력 데이터 상세 정보 추가
        if '행 ×' in output_info:
            for line in output_info.split('\n'):
                if '행 ×' in line:
                    table_data.append(['출력 형태', line.split(':')[1].strip() if ':' in line else line.strip(), ''])
                    break
        
        # 표 생성 (3열)
        rows = len(table_data)
        cols = 3
        
        table = slide.shapes.add_table(rows, cols, x, y, width, height).table
        
        # 열 너비 조정 (항목: 20%, 내용: 40%, 추가 정보: 40%)
        try:
            total_width = width
            table.columns[0].width = int(total_width * 0.2)
            table.columns[1].width = int(total_width * 0.4)
            table.columns[2].width = int(total_width * 0.4)
        except:
            pass  # 열 너비 조정 실패 시 기본값 사용
        
        # 표 스타일 설정
        for i, row_data in enumerate(table_data):
            for j, cell_data in enumerate(row_data):
                cell = table.cell(i, j)
                cell.text = str(cell_data)
                
                # 셀 텍스트 자동 줄바꿈 설정
                cell.text_frame.word_wrap = True
                cell.text_frame.margin_left = Inches(0.05)
                cell.text_frame.margin_right = Inches(0.05)
                cell.text_frame.margin_top = Inches(0.05)
                cell.text_frame.margin_bottom = Inches(0.05)
                
                # 텍스트가 너무 길면 자르기 (열에 따라 다르게)
                cell_text = str(cell_data) if cell_data else ""
                if j == 0:  # 항목 열
                    if len(cell_text) > 20:
                        cell_text = cell_text[:17] + "..."
                elif j == 1:  # 내용 열
                    if len(cell_text) > 40:
                        cell_text = cell_text[:37] + "..."
                else:  # 추가 정보 열
                    if len(cell_text) > 50:
                        cell_text = cell_text[:47] + "..."
                cell.text = cell_text
                
                # 헤더 행 스타일
                if i == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(31, 78, 121)  # 파란색 헤더
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.bold = True
                        paragraph.font.color.rgb = RGBColor(255, 255, 255)
                        paragraph.font.size = Pt(12)
                else:
                    # 짝수 행 배경색 (가독성 향상)
                    if i % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(248, 248, 248)  # 연한 회색
                    else:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(255, 255, 255)  # 흰색
                    
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(11)
                        paragraph.font.color.rgb = RGBColor(50, 50, 50)  # 진한 회색
                        paragraph.line_spacing = 1.15
        
        return table
    except Exception as e:
        print(f"표 생성 실패: {e}")
        return None

def create_comparison_chart(module, all_modules, connections):
    """데이터 분포 변화 또는 모델 성능 지표 비교 차트 생성"""
    try:
        module_type = module.get('type', 'Unknown')
        output_data = module.get('outputData')
        parameters = module.get('parameters', {})
        
        # 전처리 모듈의 경우: 입력/출력 데이터 비교
        if module_type in ['HandleMissingValues', 'EncodeCategorical', 'ScalingTransform', 'SelectData', 'TransitionData']:
            # 입력 데이터 정보
            input_info = get_input_data_info(module, all_modules, connections)
            output_info = get_output_data_info(module)
            
            # 행 수 및 열 수 추출
            input_rows = 0
            output_rows = 0
            input_cols = 0
            output_cols = 0
            
            if '행 ×' in input_info:
                try:
                    match = re.search(r'(\d+)행\s*×\s*(\d+)열', input_info)
                    if match:
                        input_rows = int(match.group(1))
                        input_cols = int(match.group(2))
                except:
                    pass
            
            if '행 ×' in output_info:
                try:
                    match = re.search(r'(\d+)행\s*×\s*(\d+)열', output_info)
                    if match:
                        output_rows = int(match.group(1))
                        output_cols = int(match.group(2))
                except:
                    pass
            
            # 행 수 비교 차트
            if input_rows > 0 and output_rows > 0:
                fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(5, 2.5))
                
                # 행 수 비교
                categories = ['입력', '출력']
                row_values = [input_rows, output_rows]
                colors = ['#4472C4', '#70AD47']
                
                bars1 = ax1.bar(categories, row_values, color=colors, alpha=0.7)
                ax1.set_ylabel('행 수', fontsize=9)
                ax1.set_title('데이터 행 수 변화', fontsize=10, fontweight='bold')
                
                for bar in bars1:
                    height = bar.get_height()
                    ax1.text(bar.get_x() + bar.get_width()/2., height,
                           f'{int(height):,}',
                           ha='center', va='bottom', fontsize=8)
                
                # 열 수 비교
                if input_cols > 0 and output_cols > 0:
                    col_values = [input_cols, output_cols]
                    bars2 = ax2.bar(categories, col_values, color=colors, alpha=0.7)
                    ax2.set_ylabel('열 수', fontsize=9)
                    ax2.set_title('데이터 열 수 변화', fontsize=10, fontweight='bold')
                    
                    for bar in bars2:
                        height = bar.get_height()
                        ax2.text(bar.get_x() + bar.get_width()/2., height,
                               f'{int(height):,}',
                               ha='center', va='bottom', fontsize=8)
                else:
                    ax2.axis('off')
                
                plt.tight_layout()
                
                # 이미지로 저장
                img_buffer = BytesIO()
                plt.savefig(img_buffer, format='png', dpi=150, bbox_inches='tight')
                plt.close()
                img_buffer.seek(0)
                
                return img_buffer
        
        # 모델 평가 모듈의 경우: 성능 지표 비교
        elif module_type == 'EvaluateModel' and output_data:
            if output_data.get('type') == 'EvaluationOutput':
                metrics = output_data.get('metrics', {})
                
                if metrics:
                    fig, ax = plt.subplots(figsize=(4, 3))
                    
                    metric_names = []
                    metric_values = []
                    
                    # 주요 지표 추출 (숫자 값만)
                    for key, value in metrics.items():
                        if isinstance(value, (int, float)) and not np.isnan(value):
                            # 지표 이름을 한국어로 변환
                            name_map = {
                                'accuracy': '정확도',
                                'precision': '정밀도',
                                'recall': '재현율',
                                'f1': 'F1 점수',
                                'f1Score': 'F1 점수',
                                'rmse': 'RMSE',
                                'mae': 'MAE',
                                'r2': 'R²',
                                'r2Score': 'R²'
                            }
                            display_name = name_map.get(key.lower(), key)
                            metric_names.append(display_name)
                            metric_values.append(value)
                    
                    if metric_names:
                        bars = ax.barh(metric_names[:5], metric_values[:5], color='#4472C4', alpha=0.7)
                        ax.set_xlabel('값', fontsize=10)
                        ax.set_title('모델 성능 지표', fontsize=11, fontweight='bold')
                        
                        # 값 표시
                        for i, (bar, val) in enumerate(zip(bars, metric_values[:5])):
                            width = bar.get_width()
                            ax.text(bar.get_width(), bar.get_y() + bar.get_height()/2.,
                                   f'{val:.3f}',
                                   ha='left', va='center', fontsize=9)
                        
                        plt.tight_layout()
                        
                        img_buffer = BytesIO()
                        plt.savefig(img_buffer, format='png', dpi=150, bbox_inches='tight')
                        plt.close()
                        img_buffer.seek(0)
                        
                        return img_buffer
        
        # 통계 모듈의 경우: 주요 통계량 시각화
        elif module_type == 'Statistics' and output_data:
            if output_data.get('type') == 'StatisticsOutput':
                stats = output_data.get('statistics', {})
                
                if stats and isinstance(stats, dict):
                    # 첫 번째 컬럼의 통계량 추출
                    first_col = list(stats.keys())[0] if stats else None
                    if first_col and isinstance(stats[first_col], dict):
                        col_stats = stats[first_col]
                        
                        fig, ax = plt.subplots(figsize=(4, 3))
                        
                        stat_names = []
                        stat_values = []
                        
                        # 주요 통계량 추출
                        for key in ['mean', 'std', 'min', 'max', 'median']:
                            if key in col_stats:
                                name_map = {
                                    'mean': '평균',
                                    'std': '표준편차',
                                    'min': '최소값',
                                    'max': '최대값',
                                    'median': '중앙값'
                                }
                                stat_names.append(name_map.get(key, key))
                                stat_values.append(col_stats[key])
                        
                        if stat_names:
                            bars = ax.bar(stat_names, stat_values, color='#4472C4', alpha=0.7)
                            ax.set_ylabel('값', fontsize=10)
                            ax.set_title(f'{first_col} 통계량', fontsize=11, fontweight='bold')
                            ax.tick_params(axis='x', rotation=45)
                            
                            # 값 표시
                            for bar in bars:
                                height = bar.get_height()
                                ax.text(bar.get_x() + bar.get_width()/2., height,
                                       f'{height:.2f}',
                                       ha='center', va='bottom', fontsize=8)
                            
                            plt.tight_layout()
                            
                            img_buffer = BytesIO()
                            plt.savefig(img_buffer, format='png', dpi=150, bbox_inches='tight')
                            plt.close()
                            img_buffer.seek(0)
                            
                            return img_buffer
        
        return None
    except Exception as e:
        print(f"차트 생성 실패: {e}")
        import traceback
        traceback.print_exc()
        return None

def get_module_characteristics(module, all_modules=None, connections=None):
    """모듈의 특성과 작업 과정을 간단히 설명"""
    module_type = module.get('type', 'Unknown')
    module_name = module.get('name', module_type)
    parameters = module.get('parameters', {})
    
    characteristics = []
    
    if module_type == 'LoadData':
        characteristics.append("특성: 외부 데이터 파일을 로드하여 분석 가능한 형태로 변환하는 데이터 입력 모듈")
        characteristics.append("작업 과정: CSV 파일을 읽어 데이터프레임으로 변환하고 다음 모듈로 전달")
    
    elif module_type == 'SelectData':
        characteristics.append("특성: 입력 데이터에서 필요한 컬럼만 선택하여 데이터 차원을 조정하는 필터링 모듈")
        characteristics.append("작업 과정: 사용자가 선택한 컬럼만 추출하여 새로운 데이터프레임 생성")
    
    elif module_type == 'HandleMissingValues':
        method = parameters.get('method', 'unknown')
        characteristics.append("특성: 결측값을 처리하여 데이터 품질을 향상시키는 전처리 모듈")
        if method == 'remove_row':
            characteristics.append("작업 과정: 결측값이 포함된 행을 완전히 제거하여 완전한 데이터만 유지")
        elif method == 'impute':
            strategy = parameters.get('strategy', 'mean')
            characteristics.append(f"작업 과정: 결측값을 {strategy} 값으로 대체하여 데이터 손실 최소화")
        elif method == 'knn':
            n_neighbors = parameters.get('n_neighbors', 5)
            characteristics.append(f"작업 과정: KNN 알고리즘을 사용하여 가장 가까운 {n_neighbors}개 이웃의 값으로 결측값 예측 및 대체")
        else:
            characteristics.append("작업 과정: 설정된 방법에 따라 결측값 처리")
    
    elif module_type == 'EncodeCategorical':
        characteristics.append("특성: 범주형 변수를 수치형으로 변환하여 머신러닝 모델에 적합하게 만드는 인코딩 모듈")
        characteristics.append("작업 과정: 선택된 범주형 컬럼을 원-핫 인코딩 또는 라벨 인코딩으로 변환하여 수치 데이터 생성")
    
    elif module_type == 'ScalingTransform':
        characteristics.append("특성: 데이터의 스케일을 정규화하여 변수 간 단위 차이를 제거하는 전처리 모듈")
        characteristics.append("작업 과정: Min-Max 스케일링 또는 Z-score 정규화를 적용하여 모든 변수를 동일한 스케일로 조정")
    
    elif module_type == 'Statistics':
        characteristics.append("특성: 데이터의 기술 통계량을 계산하여 데이터 분포와 특성을 파악하는 분석 모듈")
        characteristics.append("작업 과정: 각 컬럼별 평균, 표준편차, 최소값, 최대값, 사분위수 등 통계량 계산")
    
    elif module_type == 'SplitData':
        test_size = parameters.get('testSize', 0.2)
        characteristics.append("특성: 데이터를 훈련용과 테스트용으로 분할하여 모델 학습 및 평가를 준비하는 데이터 분할 모듈")
        characteristics.append(f"작업 과정: 데이터를 {int((1-test_size)*100)}% 훈련 데이터와 {int(test_size*100)}% 테스트 데이터로 분할")
    
    elif module_type == 'TrainModel':
        characteristics.append("특성: 학습 데이터를 사용하여 머신러닝 모델의 파라미터를 최적화하는 훈련 모듈")
        characteristics.append("작업 과정: 손실 함수를 최소화하는 방향으로 모델 파라미터를 반복적으로 업데이트하여 최적 모델 도출")
    
    elif module_type == 'ScoreModel':
        characteristics.append("특성: 훈련된 모델을 사용하여 새로운 데이터에 대한 예측값을 생성하는 예측 모듈")
        characteristics.append("작업 과정: 입력 데이터에 훈련된 모델을 적용하여 예측값을 계산하고 결과 데이터에 추가")
    
    elif module_type == 'EvaluateModel':
        characteristics.append("특성: 훈련된 모델의 성능을 평가하여 정확도, 정밀도, 재현율 등의 지표를 계산하는 평가 모듈")
        characteristics.append("작업 과정: 테스트 데이터에 대한 예측 결과와 실제 값을 비교하여 다양한 성능 지표 계산")
    
    elif module_type in ['KMeans', 'PrincipalComponentAnalysis']:
        characteristics.append("특성: 비지도 학습 알고리즘을 사용하여 데이터의 패턴이나 구조를 발견하는 분석 모듈")
        if module_type == 'KMeans':
            n_clusters = parameters.get('nClusters', 3)
            characteristics.append(f"작업 과정: K-means 클러스터링을 사용하여 데이터를 {n_clusters}개의 그룹으로 분류")
        elif module_type == 'PrincipalComponentAnalysis':
            n_components = parameters.get('nComponents', 2)
            characteristics.append(f"작업 과정: 주성분 분석을 통해 데이터를 {n_components}개의 주요 차원으로 축소")
    
    elif module_type == 'ResultModel':
        characteristics.append("특성: 통계 모델을 피팅하여 변수 간 관계를 분석하고 회귀식 또는 분류식을 도출하는 모델링 모듈")
        characteristics.append("작업 과정: 특성 변수와 목표 변수를 사용하여 통계 모델을 피팅하고 계수 및 통계량 계산")
    
    else:
        characteristics.append(f"특성: {module_name} 모듈은 데이터 처리 및 분석을 수행합니다")
        characteristics.append("작업 과정: 입력 데이터를 처리하여 결과를 출력합니다")
    
    return '\n'.join(characteristics)

def get_analysis_description(module, all_modules=None, connections=None):
    """모듈의 분석 내용 설명 생성"""
    module_type = module.get('type', 'Unknown')
    module_name = module.get('name', module_type)
    parameters = module.get('parameters', {})
    output_data = module.get('outputData')
    
    analysis_parts = []
    
    # 모듈 타입에 따른 분석 내용
    if module_type == 'LoadData':
        source = parameters.get('source', '알 수 없음')
        analysis_parts.append(f"데이터 소스: {source}")
        analysis_parts.append("CSV 파일에서 데이터를 로드하여 데이터프레임으로 변환")
        analysis_parts.append("데이터 흐름: 파일 → 데이터프레임 변환 → 다음 모듈로 전달")
    elif module_type == 'SelectData':
        selected_cols = parameters.get('columnSelections', {})
        if selected_cols:
            selected_list = [k for k, v in selected_cols.items() if v]
            analysis_parts.append(f"선택된 컬럼: {len(selected_list)}개")
            if selected_list:
                col_names = ', '.join(selected_list[:10])
                if len(selected_list) > 10:
                    col_names += f" 외 {len(selected_list) - 10}개"
                analysis_parts.append(f"  - {col_names}")
        else:
            # 모든 컬럼 선택
            if output_data and output_data.get('type') == 'DataPreview':
                columns = output_data.get('columns', [])
                analysis_parts.append(f"전체 컬럼 선택: {len(columns)}개")
        analysis_parts.append("분석 방법: 입력 데이터에서 필요한 컬럼만 필터링하여 출력")
        analysis_parts.append("데이터 흐름: 전체 데이터 → 컬럼 선택 → 선택된 데이터 출력")
    elif module_type == 'HandleMissingValues':
        method = parameters.get('method', 'unknown')
        if method == 'remove_row':
            analysis_parts.append("결측값 처리 방법: 행 제거")
            analysis_parts.append("분석 방법: 결측값이 포함된 행을 완전히 제거하여 완전한 데이터만 유지")
        elif method == 'impute':
            strategy = parameters.get('strategy', 'mean')
            analysis_parts.append(f"결측값 처리 방법: 대체 ({strategy})")
            analysis_parts.append(f"분석 방법: 결측값을 {strategy} 값으로 대체하여 데이터 손실 최소화")
        elif method == 'knn':
            n_neighbors = parameters.get('n_neighbors', 5)
            analysis_parts.append(f"결측값 처리 방법: KNN 기반 대체 (n_neighbors={n_neighbors})")
            analysis_parts.append(f"분석 방법: 가장 가까운 {n_neighbors}개 이웃의 값을 사용하여 결측값 예측")
        analysis_parts.append("데이터 흐름: 입력 데이터 → 결측값 검출 → 처리 적용 → 정제된 데이터 출력")
    elif module_type == 'Statistics':
        analysis_parts.append("분석 방법: 기술 통계량 계산")
        analysis_parts.append("- 평균, 표준편차, 최소값, 최대값, 사분위수 등")
        analysis_parts.append("- 각 컬럼별 분포 및 요약 통계 제공")
        analysis_parts.append("데이터 흐름: 입력 데이터 → 통계량 계산 → 요약 결과 출력")
    elif module_type in ['LinearRegression', 'LogisticRegression', 'PoissonRegression']:
        analysis_parts.append(f"{module_type} 모델 정의")
        analysis_parts.append("머신러닝 모델 구조 설정")
    elif module_type == 'TrainModel':
        analysis_parts.append("모델 훈련 수행")
        analysis_parts.append("학습 데이터로 모델 파라미터 최적화")
    elif module_type == 'ResultModel':
        feature_cols = parameters.get('feature_columns', [])
        label_col = parameters.get('label_column', '')
        
        # 연결된 모델 정의 모듈 찾기
        model_type = 'Unknown'
        if all_modules and connections:
            for conn in connections:
                if conn.get('to', {}).get('moduleId') == module.get('id') and \
                   conn.get('to', {}).get('portName') == 'model_in':
                    from_module_id = conn.get('from', {}).get('moduleId')
                    model_module = next((m for m in all_modules if m.get('id') == from_module_id), None)
                    if model_module:
                        if model_module.get('type') == 'StatModels':
                            model_type = model_module.get('parameters', {}).get('model', 'Unknown')
                        else:
                            model_type = model_module.get('type', 'Unknown')
                        break
        
        analysis_parts.append(f"사용 모델: {model_type}")
        analysis_parts.append(f"특성 변수: {len(feature_cols)}개 - {', '.join(feature_cols[:5])}{'...' if len(feature_cols) > 5 else ''}")
        analysis_parts.append(f"목표 변수: {label_col}")
        
        # outputData에서 모델 결과 정보 추출
        if output_data and output_data.get('type') == 'StatsModelsResultOutput':
            actual_model_type = output_data.get('modelType', model_type)
            summary = output_data.get('summary', {})
            metrics = summary.get('metrics', {})
            
            analysis_parts.append(f"적용된 모델: {actual_model_type}")
            analysis_parts.append("분석 방법:")
            if actual_model_type == 'OLS':
                analysis_parts.append("  - 최소제곱법(OLS)을 사용한 선형 회귀")
            elif actual_model_type in ['Logistic', 'Logit']:
                analysis_parts.append("  - 로지스틱 회귀를 사용한 이항 분류")
            elif actual_model_type == 'Poisson':
                analysis_parts.append("  - 포아송 회귀를 사용한 카운트 데이터 모델링")
            elif actual_model_type == 'QuasiPoisson':
                analysis_parts.append("  - 준포아송 회귀를 사용한 과분산 카운트 데이터 모델링")
            elif actual_model_type == 'NegativeBinomial':
                analysis_parts.append("  - 음이항 회귀를 사용한 과분산 카운트 데이터 모델링")
            
            # 함수식 생성
            equation = get_model_equation(output_data)
            if equation:
                analysis_parts.append(f"생성된 함수식: {equation}")
            
            # 주요 지표
            if metrics:
                key_metrics = []
                if 'R-squared' in metrics:
                    key_metrics.append(f"R² = {metrics['R-squared']}")
                if 'AIC' in metrics:
                    key_metrics.append(f"AIC = {metrics['AIC']}")
                if 'Log-Likelihood' in metrics:
                    key_metrics.append(f"Log-Likelihood = {metrics['Log-Likelihood']}")
                if key_metrics:
                    analysis_parts.append(f"주요 지표: {', '.join(key_metrics)}")
        else:
            analysis_parts.append("분석 방법: 통계 모델 피팅 및 결과 분석")
        
        analysis_parts.append("데이터 흐름: 데이터 + 모델 정의 → 모델 피팅 → 계수 및 통계량 계산 → 결과 출력")
    elif module_type == 'PredictModel':
        analysis_parts.append("분석 방법: 훈련된 모델을 사용하여 예측 수행")
        analysis_parts.append("예측 과정: 입력 데이터에 모델 함수식 적용 → 예측값 계산 → 결과 데이터에 예측값 컬럼 추가")
        analysis_parts.append("데이터 흐름: 데이터 + 훈련된 모델 → 예측 수행 → 예측값 포함 데이터 출력")
    elif module_type == 'StatModels':
        model_type = parameters.get('model', 'Unknown')
        analysis_parts.append(f"모델 정의: {model_type} 모델 구조 설정")
        analysis_parts.append("분석 방법: 통계 모델의 구조와 파라미터를 정의하여 모델 인스턴스 생성")
        analysis_parts.append("데이터 흐름: 모델 타입 선택 → 모델 정의 생성 → Result Model로 전달")
    elif module_type == 'TrainModel':
        analysis_parts.append("분석 방법: 모델 훈련 수행")
        analysis_parts.append("훈련 과정: 학습 데이터로 모델 파라미터 최적화 → 손실 함수 최소화 → 최적 파라미터 도출")
        analysis_parts.append("데이터 흐름: 데이터 + 모델 정의 → 파라미터 최적화 → 훈련된 모델 출력")
    else:
        analysis_parts.append(f"{module_type} 모듈 실행")
        if parameters:
            key_params = list(parameters.keys())[:3]
            analysis_parts.append(f"주요 파라미터: {', '.join(key_params)}")
        analysis_parts.append("데이터 흐름: 입력 데이터 처리 → 분석 수행 → 결과 출력")
    
    return '\n'.join(analysis_parts)

def create_flowchart_slide(prs, modules, connections, project_name):
    """전체 모델 흐름도를 텍스트 형태로 그리는 슬라이드 생성 (자동 줄바꿈 및 모델 정의 모듈 우선 배치)"""
    # 빈 슬라이드 생성
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃
    
    # 배경 색상 설정
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(250, 250, 252)  # 더 밝은 회색 배경
    
    # 제목 추가
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_shape.text_frame
    title_frame.text = f"전체 모델 흐름도: {project_name}"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(20)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 0, 0)  # 검정색
    title_para.alignment = PP_ALIGN.CENTER
    
    # 유효한 모듈만 필터링
    valid_modules = [m for m in modules if m.get('type') not in EXCLUDED_TYPES]
    
    if not valid_modules:
        return
    
    # 모델 정의 모듈 타입 목록 (TrainModel, TrainClusteringModel, ResultModel 앞에 와야 함)
    model_definition_types = [
        'DecisionTree', 'RandomForest', 'NeuralNetwork', 'SVM', 'LDA', 'NaiveBayes', 'KNN',
        'LinearRegression', 'LogisticRegression', 'PoissonRegression', 'NegativeBinomialRegression',
        'KMeans', 'PCA'
    ]
    
    # 모듈 순서 결정 (위상 정렬 + 모델 정의 모듈 우선 배치)
    def get_module_level(module_id, visited=None):
        if visited is None:
            visited = set()
        if module_id in visited:
            return 0
        visited.add(module_id)
        
        # 이 모듈로 들어오는 연결 찾기
        incoming = [c for c in connections if c.get('to', {}).get('moduleId') == module_id]
        if not incoming:
            return 0
        
        # 최대 깊이 계산
        max_depth = 0
        for conn in incoming:
            from_id = conn.get('from', {}).get('moduleId')
            depth = get_module_level(from_id, visited.copy())
            max_depth = max(max_depth, depth)
        
        return max_depth + 1
    
    # 모듈들을 레벨별로 그룹화
    modules_by_level = {}
    for module in valid_modules:
        level = get_module_level(module.get('id'))
        if level not in modules_by_level:
            modules_by_level[level] = []
        modules_by_level[level].append(module)
    
    # 레벨 순서대로 모듈 이름 수집 (모델 정의 모듈 우선 배치)
    module_sequence = []
    for level in sorted(modules_by_level.keys()):
        level_modules = modules_by_level[level]
        
        # 모델 정의 모듈과 일반 모듈 분리
        model_def_modules = [m for m in level_modules if m.get('type') in model_definition_types]
        other_modules = [m for m in level_modules if m.get('type') not in model_definition_types]
        
        # TrainModel, TrainClusteringModel, ResultModel 찾기
        train_modules = [m for m in other_modules if m.get('type') in ['TrainModel', 'TrainClusteringModel', 'ResultModel']]
        other_other_modules = [m for m in other_modules if m.get('type') not in ['TrainModel', 'TrainClusteringModel', 'ResultModel']]
        
        # 순서: 모델 정의 모듈 -> TrainModel/TrainClusteringModel/ResultModel -> 기타 모듈
        # 같은 그룹 내에서는 이름 순으로 정렬
        model_def_modules.sort(key=lambda m: m.get('name', m.get('type', '')))
        train_modules.sort(key=lambda m: m.get('name', m.get('type', '')))
        other_other_modules.sort(key=lambda m: m.get('name', m.get('type', '')))
        
        module_sequence.extend(model_def_modules)
        module_sequence.extend(train_modules)
        module_sequence.extend(other_other_modules)
    
    # 텍스트 형태로 흐름도 생성: "Module1 -> Module2 -> Module3"
    module_names = [m.get('name', m.get('type', 'Unknown')) for m in module_sequence]
    flow_text = " -> ".join(module_names)
    
    # 텍스트 박스 생성 (중앙 배치, 넓게, 여러 줄 가능)
    text_shape = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(4.5))
    text_frame = text_shape.text_frame
    text_frame.word_wrap = True  # 자동 줄바꿈 활성화 (슬라이드를 넘어가면 아래로 이동)
    text_frame.margin_left = Inches(0.2)
    text_frame.margin_right = Inches(0.2)
    text_frame.margin_top = Inches(0.2)
    text_frame.margin_bottom = Inches(0.2)
    
    # 전체 텍스트를 하나의 단락으로 설정
    para = text_frame.paragraphs[0]
    para.text = flow_text
    para.font.size = Pt(17)
    para.font.bold = True
    para.font.color.rgb = RGBColor(31, 78, 121)  # 파란색
    para.alignment = PP_ALIGN.CENTER
    para.space_after = Pt(0)
    para.space_before = Pt(0)
    
    # 텍스트 프레임 자동 크기 조정 비활성화
    text_frame.auto_size = MSO_AUTO_SIZE.NONE

def truncate_text_to_fit(text, max_lines=10, max_chars_per_line=80):
    """텍스트를 지정된 줄 수와 문자 수로 자르기"""
    lines = text.split('\n')
    truncated_lines = []
    total_chars = 0
    
    for line in lines[:max_lines]:
        if len(line) > max_chars_per_line:
            # 긴 줄을 여러 줄로 분할
            while len(line) > max_chars_per_line:
                truncated_lines.append(line[:max_chars_per_line])
                line = line[max_chars_per_line:]
            if line:
                truncated_lines.append(line)
        else:
            truncated_lines.append(line)
        
        if len(truncated_lines) >= max_lines:
            break
    
    result = '\n'.join(truncated_lines)
    if len(lines) > max_lines or any(len(line) > max_chars_per_line for line in lines):
        result += f"\n... (내용 일부 생략)"
    
    return result

def create_module_slide(prs, module, all_modules, connections, module_index, total_modules):
    """하나의 모듈에 대한 슬라이드 생성 (개선된 버전: 헤드라인, 표, 시각적 요소 포함)"""
    module_type = module.get('type', 'Unknown')
    module_name = module.get('name', module_type)
    
    # 슬라이드 추가 (빈 레이아웃 사용)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃
    
    # 배경 색상 설정 (더 밝은 회색으로 변경)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(250, 250, 252)  # 더 밝은 회색 배경
    
    # 1. 분석적 통찰 헤드라인 (상단 중앙) - 배경 박스 추가
    headline = generate_insight_headline(module, all_modules, connections)
    
    # 모듈 이름을 헤드라인 뒤에 추가
    module_name = module.get('name', module.get('type', 'Unknown'))
    if headline and not headline.endswith(module_name):
        # 헤드라인과 모듈 이름을 구분하여 표시
        full_headline = f"{headline} - {module_name}"
    else:
        full_headline = headline if headline else f"{module_name} 분석 결과"
    
    # 헤드라인 배경 박스 (그림자 효과를 위한 추가 레이어)
    headline_bg_shadow = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.32), Inches(0.17), Inches(9.4), Inches(0.7)
    )
    headline_bg_shadow.fill.solid()
    headline_bg_shadow.fill.fore_color.rgb = RGBColor(20, 50, 80)  # 그림자 색상
    headline_bg_shadow.line.fill.background()
    
    headline_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.7)
    )
    headline_bg.fill.solid()
    headline_bg.fill.fore_color.rgb = RGBColor(31, 78, 121)  # 진한 파란색 배경
    headline_bg.line.fill.background()  # 테두리 없음
    
    headline_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    headline_frame = headline_shape.text_frame
    headline_frame.text = full_headline
    headline_frame.word_wrap = True
    headline_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    headline_para = headline_frame.paragraphs[0]
    headline_para.font.size = Pt(20)
    headline_para.font.bold = True
    headline_para.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 텍스트
    headline_para.alignment = PP_ALIGN.CENTER
    headline_para.space_after = Pt(0)
    
    # 2. 모듈 번호 및 이름 (좌측 상단, 작게) - 제거하고 공간 확보
    # 모듈 정보는 헤드라인에 포함되므로 별도 표시 제거
    
    # 3. 입력/출력 데이터 요약 표 (좌측 상단) - 위치 조정 및 배경 추가
    table_y = Inches(1.0)
    table_width = Inches(4.5)
    table_height = Inches(2.0)
    
    # 표 배경 박스
    table_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.25), Inches(0.95), Inches(4.6), Inches(2.1)
    )
    table_bg.fill.solid()
    table_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)  # 흰색 배경
    table_bg.line.color.rgb = RGBColor(200, 200, 200)  # 연한 회색 테두리
    table_bg.line.width = Pt(1)
    
    create_data_summary_table(slide, module, all_modules, connections, 
                             Inches(0.35), table_y, table_width, table_height)
    
    # 4. 시각적 요소 (차트 또는 그래프) - 우측 상단 - 위치 조정
    chart_img = create_comparison_chart(module, all_modules, connections)
    if chart_img:
        try:
            # 임시 파일로 저장
            with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp_file:
                tmp_file.write(chart_img.read())
                tmp_path = tmp_file.name
            
            # 차트 배경 박스 추가
            chart_bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(5.15), Inches(0.95), Inches(4.4), Inches(2.1)
            )
            chart_bg.fill.solid()
            chart_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)  # 흰색 배경
            chart_bg.line.color.rgb = RGBColor(200, 200, 200)  # 연한 회색 테두리
            chart_bg.line.width = Pt(1)
            
            # 이미지 삽입 (위치 조정)
            slide.shapes.add_picture(tmp_path, Inches(5.25), Inches(1.0), width=Inches(4.2), height=Inches(2.0))
            
            # 임시 파일 삭제
            try:
                os.unlink(tmp_path)
            except:
                pass
        except Exception as e:
            print(f"차트 이미지 삽입 실패: {e}")
    
    # 5. 모듈 특성 및 작업 과정 (좌측 하단) - 배경 박스 추가
    characteristics = get_module_characteristics(module, all_modules, connections)
    
    # 배경 박스 추가
    char_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.25), Inches(3.05), Inches(4.6), Inches(2.1)
    )
    char_bg.fill.solid()
    char_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)  # 흰색 배경
    char_bg.line.color.rgb = RGBColor(200, 200, 200)  # 연한 회색 테두리
    char_bg.line.width = Pt(1)
    
    char_shape = slide.shapes.add_textbox(Inches(0.35), Inches(3.15), Inches(4.4), Inches(1.9))
    char_frame = char_shape.text_frame
    char_frame.word_wrap = True
    char_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    char_frame.margin_left = Inches(0.15)
    char_frame.margin_right = Inches(0.15)
    char_frame.margin_top = Inches(0.1)
    char_frame.margin_bottom = Inches(0.1)
    
    p = char_frame.paragraphs[0]
    p.text = "🔧 모듈 특성 및 작업 과정"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(31, 78, 121)  # 파란색 제목
    p.space_after = Pt(4)
    
    # 텍스트를 줄 단위로 추가하고 길면 자동 줄바꿈
    char_lines = characteristics.split('\n')
    displayed_lines = 0
    max_lines = 8  # 최대 8줄
    
    for line in char_lines:
        if line.strip() and displayed_lines < max_lines:
            p = char_frame.add_paragraph()
            # 65자 초과 시 자동 줄바꿈
            display_text = line[:65] if len(line) <= 65 else line[:62] + "..."
            p.text = display_text
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(50, 50, 50)  # 진한 회색 텍스트
            p.space_after = Pt(2)
            p.line_spacing = 1.2
            displayed_lines += 1
    
    # 텍스트가 너무 많으면 생략 표시
    if len(char_lines) > max_lines:
        p = char_frame.add_paragraph()
        p.text = f"... (총 {len(char_lines)}줄 중 {max_lines}줄만 표시)"
        p.font.size = Pt(9)
        p.font.italic = True
        p.font.color.rgb = RGBColor(150, 150, 150)
    
    # 6. 상세 분석 내용 (우측 하단) - 배경 박스 추가
    analysis_desc = get_analysis_description(module, all_modules, connections)
    
    # 배경 박스 추가
    analysis_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.15), Inches(3.05), Inches(4.4), Inches(2.1)
    )
    analysis_bg.fill.solid()
    analysis_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)  # 흰색 배경
    analysis_bg.line.color.rgb = RGBColor(200, 200, 200)  # 연한 회색 테두리
    analysis_bg.line.width = Pt(1)
    
    analysis_shape = slide.shapes.add_textbox(Inches(5.25), Inches(3.15), Inches(4.2), Inches(1.9))
    analysis_frame = analysis_shape.text_frame
    analysis_frame.word_wrap = True
    analysis_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    analysis_frame.margin_left = Inches(0.15)
    analysis_frame.margin_right = Inches(0.15)
    analysis_frame.margin_top = Inches(0.1)
    analysis_frame.margin_bottom = Inches(0.1)
    
    p = analysis_frame.paragraphs[0]
    p.text = "📊 상세 분석 내용"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(31, 78, 121)  # 파란색 제목
    p.space_after = Pt(4)
    
    # 분석 내용을 여러 줄로 분리 (최대 10줄)
    analysis_lines = [line for line in analysis_desc.split('\n') if line.strip()]
    displayed_lines = 0
    max_lines = 10  # 최대 10줄
    
    for line in analysis_lines:
        if displayed_lines < max_lines:
            p = analysis_frame.add_paragraph()
            # 65자 초과 시 자동 줄바꿈
            display_text = line[:65] if len(line) <= 65 else line[:62] + "..."
            p.text = display_text
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(50, 50, 50)  # 진한 회색 텍스트
            p.space_after = Pt(2)
            p.line_spacing = 1.2
            displayed_lines += 1
    
    # 텍스트가 너무 많으면 생략 표시
    if len(analysis_lines) > max_lines:
        p = analysis_frame.add_paragraph()
        p.text = f"... (총 {len(analysis_lines)}줄 중 {max_lines}줄만 표시)"
        p.font.size = Pt(9)
        p.font.italic = True
        p.font.color.rgb = RGBColor(150, 150, 150)
    
    # 7. 주요 파라미터 (표 아래, 좌측) - 제거하고 공간 확보
    # 파라미터는 표에 포함되므로 별도 섹션 제거

def create_single_ppt(project_data, output_path):
    """하나의 PPT 파일에 모든 모듈 슬라이드 생성"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    modules = project_data.get('modules', [])
    connections = project_data.get('connections', [])
    project_name = project_data.get('projectName', 'Untitled Project')
    
    # TextBox, GroupBox 같은 Shape 타입은 제외
    valid_modules = [m for m in modules if m.get('type') not in EXCLUDED_TYPES]
    
    if not valid_modules:
        print("유효한 모듈이 없습니다.")
        return None
    
    # 제목 슬라이드
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    # 제목 슬라이드 배경 설정
    background = title_slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(31, 78, 121)  # 진한 파란색 배경
    
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    
    title.text = project_name
    # 제목 텍스트 색상 설정 (흰색)
    for paragraph in title.text_frame.paragraphs:
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.font.size = Pt(44)
        paragraph.font.bold = True
    
    subtitle.text = f"총 {len(valid_modules)}개 모듈"
    # 부제목 텍스트 색상 설정 (연한 흰색)
    for paragraph in subtitle.text_frame.paragraphs:
        paragraph.font.color.rgb = RGBColor(230, 230, 230)
        paragraph.font.size = Pt(24)
    
    # 전체 흐름도 슬라이드 추가
    print("전체 흐름도 슬라이드 생성 중...")
    create_flowchart_slide(prs, modules, connections, project_name)
    
    # 각 모듈에 대한 슬라이드 생성
    print(f"총 {len(valid_modules)}개의 모듈에 대해 슬라이드를 생성합니다...")
    for i, module in enumerate(valid_modules, 1):
        try:
            module_name = module.get('name', module.get('type'))
            print(f"[{i}/{len(valid_modules)}] {module_name} 슬라이드 생성 중...")
            create_module_slide(prs, module, modules, connections, i, len(valid_modules))
        except Exception as e:
            print(f"모듈 {module.get('name', module.get('type'))} 슬라이드 생성 실패: {e}")
            import traceback
            traceback.print_exc()
    
    # 파일 저장
    prs.save(output_path)
    print(f"\nPPT 파일 생성 완료: {output_path}")
    print(f"다운로드 폴더에 저장되었습니다: {os.path.dirname(output_path)}")
    
    return output_path

def process_project_data(project_data, output_dir=None):
    """프로젝트 데이터를 읽어서 하나의 PPT 파일 생성"""
    # 출력 디렉토리 설정 (다운로드 폴더 또는 지정된 폴더)
    if output_dir is None:
        output_dir = get_download_folder()
        print(f"다운로드 폴더에 저장: {output_dir}")
    else:
        os.makedirs(output_dir, exist_ok=True)
        print(f"지정된 폴더에 저장: {output_dir}")
    
    # 다운로드 폴더가 존재하는지 확인
    if not os.path.exists(output_dir):
        print(f"경고: 다운로드 폴더가 존재하지 않습니다. 생성합니다: {output_dir}")
        os.makedirs(output_dir, exist_ok=True)
    
    project_name = project_data.get('projectName', 'Untitled_Project')
    safe_name = sanitize_filename(project_name)
    filename = f"{safe_name}_모듈분석.pptx"
    output_path = os.path.join(output_dir, filename)
    print(f"저장 경로: {output_path}")
    
    # PPT 파일 생성
    result_path = create_single_ppt(project_data, output_path)
    
    if result_path:
        return [{
            'filename': filename,
            'filepath': result_path,
            'module_count': len([m for m in project_data.get('modules', []) if m.get('type') not in EXCLUDED_TYPES])
        }]
    else:
        return []

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("사용법: python generate_module_ppts.py <project_json_file> [output_dir]")
        print("  output_dir을 지정하지 않으면 다운로드 폴더에 저장됩니다.")
        sys.exit(1)
    
    input_file = sys.argv[1]
    output_dir = sys.argv[2] if len(sys.argv) > 2 else None
    
    if not os.path.exists(input_file):
        print(f"파일을 찾을 수 없습니다: {input_file}")
        sys.exit(1)
    
    try:
        with open(input_file, 'r', encoding='utf-8') as f:
            project_data = json.load(f)
        
        generated_files = process_project_data(project_data, output_dir)
        
        if generated_files:
            print(f"\n생성 완료: {generated_files[0]['filepath']}")
            print(f"모듈 수: {generated_files[0]['module_count']}개")
        else:
            print("\nPPT 파일 생성에 실패했습니다.")
        
    except Exception as e:
        print(f"오류 발생: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)
